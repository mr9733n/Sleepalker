from pptx import Presentation

prs = Presentation()

# Slide 1 Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Sleepwalker — Pitch Deck"
subtitle.text = "Комедийный социальный стелс"

# Slide 2 Concept
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "Концепция"
body.text = ("Главный герой страдает лунатизмом и просыпается в случайных местах без одежды.\n"
             "Цель: найти одежду и избежать социального позора, используя стелс, юмор и взаимодействие с NPC.")

# Slide 3 Key Features
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "Ключевые особенности"
body.text = ("• Комедийный социальный стелс\n"
             "• Процедурная генерация ситуаций\n"
             "• Реакции NPC на основе AI\n"
             "• Импровизация одежды и действий\n"
             "• Множество уникальных сценариев пробуждения")

# Slide 4 Gameplay Loop
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "Игровой цикл"
body.text = ("1. ГГ засыпает\n"
             "2. Случайное пробуждение\n"
             "3. Поиск одежды / сокрытие\n"
             "4. Стелс и взаимодействие\n"
             "5. Возвращение в безопасную зону\n"
             "6. Прогрессия персонажа")

# Slide 5 Scenarios
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "Примеры ситуаций"
body.text = ("• Лифт в отеле\n"
             "• Торговый центр (витрина манекенов)\n"
             "• Пляж под пирсом\n"
             "• Фонтан на площади\n"
             "• Крыша многоэтажки\n"
             "• Багажная полка поезда")

# Slide 6 NPC
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "NPC и их поведение"
body.text = ("• Охранники, туристы, работники, дети, полиция\n"
             "• Реакции: страх, смех, запись на телефон, помощь\n"
             "• AI: поле зрения, эмпатия, тревожность, маршруты")

# Slide 7 Improvisation
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "Импровизация"
body.text = ("Предметы: полотенца, газеты, пакеты, шторы.\n"
             "Игрок придумывает способы скрыться, отвлечь NPC или добыть одежду.")

# Slide 8 Target Audience
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "Целевая аудитория"
body.text = ("• Любители стелса\n"
             "• Игроки, ценящие юмор\n"
             "• Фанаты игр Untitled Goose Game, Hitman (социальный стелс)\n"
             "• Возраст 16+")

# Slide 9 Production
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "Производство"
body.text = ("Мини-команда:\n"
             "• 1 программист\n"
             "• 1 геймдизайнер\n"
             "• 1 3D-художник\n"
             "• 1 дизайнер уровней\n\n"
             "Срок на MVP: 3–5 месяцев")

# Slide 10 Summary
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "Резюме"
body.text = ("Игра предлагает уникальный комедийный подход к стелсу.\n"
             "Социальные ситуации, процедурность и юмор делают проект ярким и заметным.\n"
             "Идеален для инди-прототипа с перспективой расширения.")

file_path = "/mnt/data/Sleepwalker_Pitch_Deck.pptx"
prs.save(file_path)

file_path
